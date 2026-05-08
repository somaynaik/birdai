from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(51, 152, 219)  # Blue background
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(32)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(title, content_list):
    """Add a slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(52, 73, 94)
        p.space_before = Pt(12)
        p.space_after = Pt(12)
        p.level = 0

# Slide 1: Title
add_title_slide("🐦 Bird Species Classifier", "Complete Project Logic & Architecture")

# Slide 2: Project Overview
add_content_slide("Project Overview", [
    "✓ AI-powered application for identifying bird species",
    "✓ Accepts TWO inputs: Bird images and bird sounds",
    "✓ Uses Deep Learning models to classify species",
    "✓ Available in TWO formats: Desktop (Tkinter) & Web (Flask)",
    "✓ Real-time predictions with high accuracy"
])

# Slide 3: Key Features
add_content_slide("Key Features", [
    "📷 Image Classification - Identify birds from photos",
    "🎵 Audio Classification - Identify birds from sound recordings",
    "🔄 Real-time Processing - Instant predictions",
    "🎨 User-friendly Interface - Easy to use for everyone",
    "⚡ Multi-threaded Loading - Smooth user experience",
    "☁️ Uses Pre-trained Models - No training needed"
])

# Slide 4: Architecture Overview
add_content_slide("System Architecture", [
    "INPUT LAYER",
    "  • Image File (.jpg, .png, .bmp)",
    "  • Audio File (.wav, .flac, .mp3)",
    "",
    "PROCESSING LAYER",
    "  • Image Processor / Audio Feature Extractor",
    "",
    "AI MODEL LAYER",
    "  • Vision Transformer (Image)",
    "  • wav2vec2 Model (Audio)",
    "",
    "OUTPUT LAYER",
    "  • Bird Species Prediction"
])

# Slide 5: Image Classification Logic
add_content_slide("Image Classification Pipeline", [
    "1. USER UPLOADS IMAGE",
    "   → Select bird photo file",
    "",
    "2. IMAGE PREPROCESSING",
    "   → Load image using PIL",
    "   → Convert to RGB format",
    "   → Resize & normalize using AutoImageProcessor",
    "",
    "3. MODEL INFERENCE",
    "   → Pass processed image to ViT/CNN model",
    "   → Extract feature embeddings",
    "",
    "4. PREDICTION",
    "   → Get logits (raw scores) from model",
    "   → Apply argmax to get highest probability class",
    "   → Return bird species name"
])

# Slide 6: Audio Classification Logic
add_content_slide("Audio Classification Pipeline", [
    "1. USER UPLOADS AUDIO FILE",
    "   → Select .wav, .flac, or .mp3 file",
    "",
    "2. AUDIO PREPROCESSING (librosa)",
    "   → Load audio file",
    "   → Resample to 16kHz (standard rate)",
    "   → Convert to MONO (single channel)",
    "   → Extract sound features",
    "",
    "3. FEATURE EXTRACTION (wav2vec2)",
    "   → Process audio through feature extractor",
    "   → Convert to numerical embeddings",
    "",
    "4. MODEL INFERENCE",
    "   → Pass embeddings to audio classification model",
    "   → Get prediction logits",
    "",
    "5. OUTPUT SPECIES",
    "   → argmax() gets highest probability"
])

# Slide 7: Models Deep Dive
add_content_slide("AI Models Used", [
    "IMAGE MODEL: chriamue/bird-species-classifier",
    "  • Type: Vision Transformer (ViT)",
    "  • Uses: Self-attention mechanism for images",
    "  • Divides image into patches → Analyzes relationships",
    "",
    "AUDIO MODEL: greenarcade/wav2vec2-vd-bird-sound",
    "  • Type: wav2vec2 (Self-supervised learning)",
    "  • Uses: Transformer architecture on audio",
    "  • Pre-trained on large unlabeled audio → Fine-tuned for birds"
])

# Slide 8: Technologies & Libraries
add_content_slide("Tech Stack", [
    "🐍 BACKEND",
    "  • Python 3.13",
    "  • PyTorch - Deep learning framework",
    "  • Transformers - Pre-trained models library",
    "  • librosa - Audio signal processing",
    "",
    "🖥️ FRONTEND",
    "  • Tkinter - Desktop GUI (main.py)",
    "  • Flask - Web framework (app.py)",
    "  • HTML/CSS/JavaScript - Web interface"
])

# Slide 9: Desktop App (Tkinter)
add_content_slide("Desktop Version (main.py)", [
    "UI Components:",
    "  • Upload Image Button → Opens file dialog",
    "  • Upload Audio Button → Opens file dialog",
    "  • Preview Area → Shows selected file",
    "  • Result Area → Shows prediction + species name",
    "  • Next Bird Button → Reset for new prediction",
    "",
    "Workflow:",
    "  1. App loads models on startup (background thread)",
    "  2. User selects image/audio",
    "  3. Processing happens in background",
    "  4. Result displays in GUI",
    "  5. User can predict again"
])

# Slide 10: Web App (Flask)
add_content_slide("Web Version (app.py)", [
    "Backend Endpoints:",
    "  • / → Serves HTML interface",
    "  • /status → Check if models are loaded",
    "  • /classify-image → Process image",
    "  • /classify-audio → Process audio",
    "",
    "Frontend (JavaScript):",
    "  • Polls /status until models load",
    "  • Sends files to Flask via POST requests",
    "  • Displays results in browser",
    "",
    "Run: python app.py",
    "Access: http://localhost:5000"
])

# Slide 11: Data Flow Diagram
add_content_slide("Complete Data Flow", [
    "USER INPUT (Image/Audio)",
    "         ↓",
    "FILE UPLOAD (Form/Dialog)",
    "         ↓",
    "PREPROCESSING (PIL/librosa)",
    "         ↓",
    "FEATURE EXTRACTION (AutoProcessor/AutoExtractor)",
    "         ↓",
    "DEEP LEARNING MODEL (ViT/wav2vec2)",
    "         ↓",
    "OUTPUT PREDICTIONS (Logits)",
    "         ↓",
    "ARGMAX (Get highest probability)",
    "         ↓",
    "SPECIES NAME (Display to user)"
])

# Slide 12: Key Concepts Explained
add_content_slide("Key AI Concepts", [
    "TRANSFORMER: Neural network architecture using attention",
    "  → Analyzes relationships between all input elements",
    "",
    "VISION TRANSFORMER (ViT): Applies transformers to images",
    "  → Treats image patches like words in sentences",
    "",
    "WAV2VEC2: Self-supervised audio representation learning",
    "  → Learns from unlabeled audio data",
    "  → Converts sound to meaningful features",
    "",
    "INFERENCE: Using trained model to make predictions",
    "  → No training happening, just prediction",
    "",
    "ARGMAX: Gets index of highest value in array",
    "  → Converts [0.1, 0.8, 0.1] → index 1 (bird species)"
])

# Slide 13: Performance Metrics
add_content_slide("Model Performance", [
    "IMAGE CLASSIFICATION:",
    "  • Models: Trained on thousands of bird images",
    "  • Accuracy: High confidence for clear photos",
    "  • Speed: ~1-2 seconds per image",
    "",
    "AUDIO CLASSIFICATION:",
    "  • Models: Fine-tuned on bird sound datasets",
    "  • Accuracy: Excellent for 5-30 second clips",
    "  • Speed: ~1-3 seconds per audio file",
    "",
    "NOTE: Quality of input affects accuracy",
    "  → Clear images = Better results",
    "  → Clean audio = Better results"
])

# Slide 14: Error Handling
add_content_slide("Error Handling & Edge Cases", [
    "Image Issues:",
    "  • Invalid format → Error message shown",
    "  • Corrupted file → Exception caught and reported",
    "",
    "Audio Issues:",
    "  • Unsupported format → librosa fails gracefully",
    "  • Empty audio → Validation check prevents crash",
    "  • Very short clips → Model may give low confidence",
    "",
    "Model Loading:",
    "  • Networks offline → Downloads cached models",
    "  • Large files → Streamed from Hugging Face",
    "  • Loading status → Status bar shows progress"
])

# Slide 15: Future Improvements
add_content_slide("Potential Enhancements", [
    "🔧 TECHNICAL",
    "  • Add confidence scores to predictions",
    "  • Support for batch processing (multiple files)",
    "  • Model quantization for faster inference",
    "  • GPU acceleration for speedup",
    "",
    "📊 FEATURES",
    "  • Database to store prediction history",
    "  • User authentication for web version",
    "  • Mobile app (React Native / Flutter)",
    "  • Real-time bird detection from camera/mic",
    "  • Multi-species probability rankings"
])

# Slide 16: Summary
add_content_slide("Summary", [
    "✅ TWO INPUT TYPES: Images & Audio recordings",
    "✅ TWO INTERFACES: Desktop (Tkinter) & Web (Flask)",
    "✅ STATE-OF-ART MODELS: Vision Transformer + wav2vec2",
    "✅ AUTOMATIC PREPROCESSING: Handles all formats",
    "✅ REAL-TIME PREDICTIONS: Instant results",
    "✅ USER-FRIENDLY: No ML knowledge needed",
    "✅ SCALABLE: Can serve multiple users (web version)"
])

# Slide 17: Questions
add_title_slide("Questions?", "🐦 Bird Species Classifier 🐦")

# Save presentation
prs.save('Bird_Species_Classifier_PPT.pptx')
print("✅ PowerPoint created: Bird_Species_Classifier_PPT.pptx")
